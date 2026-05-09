"""Slide parsing and simple retrieval utilities."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re

from pypdf import PdfReader

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency may fail in some envs
    Presentation = None  # type: ignore[assignment]


@dataclass(frozen=True)
class SlideChunk:
    page: int
    text: str


def _clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    return text


def parse_pdf(path: str | Path) -> list[SlideChunk]:
    reader = PdfReader(str(path))
    chunks: list[SlideChunk] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = _clean_text(page.extract_text() or "")
        if text:
            chunks.append(SlideChunk(page=idx, text=text))
    if not chunks:
        raise ValueError(
            "No readable text was found in the PDF. Try exporting slides with selectable text, "
            "or add OCR as a bonus feature."
        )
    return chunks


def parse_pptx(path: str | Path) -> list[SlideChunk]:
    if Presentation is None:
        raise RuntimeError("python-pptx is not available in this environment.")
    prs = Presentation(str(path))
    chunks: list[SlideChunk] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        text = _clean_text("\n".join(parts))
        if text:
            chunks.append(SlideChunk(page=idx, text=text))
    if not chunks:
        raise ValueError("No readable text was found in the PPTX.")
    return chunks


def parse_slides(path: str | Path) -> list[SlideChunk]:
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix == ".pptx":
        return parse_pptx(path)
    raise ValueError("Unsupported file type. Please upload a PDF; PPTX is supported as a bonus.")


def chunks_to_context(chunks: list[SlideChunk], max_chars: int = 12000) -> str:
    """Return slide text with references, trimmed to fit a local model context."""
    rendered: list[str] = []
    total = 0
    for chunk in chunks:
        piece = f"[Slide {chunk.page}] {chunk.text}"
        if total + len(piece) > max_chars:
            remaining = max_chars - total
            if remaining > 100:
                rendered.append(piece[:remaining] + " ...")
            break
        rendered.append(piece)
        total += len(piece)
    return "\n\n".join(rendered)


def rough_summary(chunks: list[SlideChunk], max_slides: int = 6) -> str:
    """A deterministic fallback summary used only when explicitly enabled."""
    selected = chunks[:max_slides]
    bullets = [f"- Slide {c.page}: {c.text[:180]}" for c in selected]
    return "\n".join(bullets)


def retrieve_relevant_chunks(chunks: list[SlideChunk], query: str, k: int = 5) -> list[SlideChunk]:
    """Simple lexical retriever for lightweight RAG without vector DB overhead."""
    terms = {t.lower() for t in re.findall(r"[A-Za-z][A-Za-z0-9_-]+", query)}
    scored: list[tuple[int, SlideChunk]] = []
    for chunk in chunks:
        words = set(re.findall(r"[A-Za-z][A-Za-z0-9_-]+", chunk.text.lower()))
        score = len(terms & words)
        scored.append((score, chunk))
    scored.sort(key=lambda item: item[0], reverse=True)
    return [chunk for score, chunk in scored[:k] if score > 0] or chunks[:k]
